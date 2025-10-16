# document_processor.py (modificado)
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
import comtypes.client
import os
import tempfile
import time
from abc import ABC, abstractmethod

class BaseProcessor(ABC):
    def __init__(self, template_path):
        self.template_path = template_path
        self.doc = None
        self.temp_files = []  # Lista para rastrear archivos temporales

    @abstractmethod
    def process(self, data_map: dict, font_map: dict):
        pass

    @abstractmethod
    def save_as_pdf(self, output_path: str):
        pass

    def _get_clean_filename(self, text: str) -> str:
        return "".join(c for c in str(text) if c.isalnum() or c in (' ', '-', '_')).rstrip()

    def _cleanup_temp_files(self):
        """Limpia todos los archivos temporales creados"""
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                print(f"⚠️ No se pudo eliminar archivo temporal {temp_file}: {e}")
        self.temp_files = []

class PdfProcessor(BaseProcessor):
    def __init__(self, template_path):
        super().__init__(template_path)
        self.doc = fitz.open(template_path)

    def _get_text_fit_info(self, rect, text, font_info, align_center=True):
        """
        Calcula posición y tamaño para texto
        Args:
            align_center: Si es True, centra el texto. Si es False, alinea a la izquierda
        """
        font_name = self._get_pdf_font(font_info['family'], font_info['bold'])
        max_width = rect.width * 0.98
        current_font_size = font_info['size']
        tw_initial = fitz.get_text_length(text, fontname=font_name, fontsize=current_font_size)

        final_font_size = current_font_size
        
        if tw_initial > max_width:
            while final_font_size > 6: 
                tw = fitz.get_text_length(text, fontname=font_name, fontsize=final_font_size)
                if tw <= max_width:
                    break
                final_font_size -= 1
        
        tw = fitz.get_text_length(text, fontname=font_name, fontsize=final_font_size)
        
        # Diferente alineación según el parámetro
        if align_center:
            x_insert = rect.x0 + (rect.width - tw) / 2  # Centrado
        else:
            x_insert = rect.x0 + 5  # Alineado a la izquierda con pequeño margen
        
        font_ascender = 0.8
        y_insert = rect.y0 + (rect.height - (final_font_size * font_ascender)) / 2 + final_font_size

        return x_insert, y_insert, final_font_size, font_name

    def process(self, data_map: dict, font_map: dict):
        self.doc = fitz.open(self.template_path)
        for page in self.doc:
            # Buscar TODOS los placeholders, incluyendo en cuadros de texto
            all_text_instances = {}
            
            for placeholder in data_map.keys():
                # Buscar en texto normal
                text_instances = page.search_for(placeholder)
                if text_instances:
                    all_text_instances[placeholder] = text_instances
                
                # Buscar en cuadros de texto (widget annotations)
                try:
                    for widget in page.widgets():
                        if widget.field_type == fitz.PDF_WIDGET_TYPE_TEXT:
                            if placeholder in widget.field_value or placeholder in widget.field_name:
                                # Crear un rectángulo aproximado para el cuadro de texto
                                rect = widget.rect
                                all_text_instances.setdefault(placeholder, []).append(rect)
                except Exception:
                    pass
            
            # Procesar cada placeholder encontrado
            for placeholder, instances in all_text_instances.items():
                value = data_map[placeholder]
                font_info = font_map.get(placeholder, {'family': 'Arial', 'size': 12, 'bold': False, 'color': (0, 0, 0)})
                
                # Determinar alineación: centrado para TEXT_1 y TEXT_2, izquierda para FOLIO
                is_folio = placeholder == "{{FOLIO}}"
                align_center = not is_folio  # Centrado para todo excepto FOLIO
                
                for inst in instances:
                    try:
                        page.add_redact_annot(inst)
                        page.apply_redactions()
                        
                        x_insert, y_insert, final_size, font_name = self._get_text_fit_info(
                            inst, value, font_info, align_center=align_center
                        )
                        
                        # Obtener color del font_info
                        color = self._parse_color(font_info.get('color', (0, 0, 0)))
                        
                        page.insert_text(
                            (x_insert, y_insert),
                            value,
                            fontname=font_name,
                            fontsize=final_size,
                            color=color
                        )
                    except Exception as e:
                        print(f"⚠️ Error procesando {placeholder}: {e}")
                        continue

    def _parse_color(self, color):
        """Convierte el color a formato compatible con PyMuPDF"""
        if isinstance(color, str) and color.startswith('#'):
            # Convertir hex a RGB (valores 0-255)
            color = color.lstrip('#')
            r = int(color[0:2], 16)
            g = int(color[2:4], 16)
            b = int(color[4:6], 16)
            return (r/255, g/255, b/255)  # Convertir a valores 0-1 para PyMuPDF
        elif isinstance(color, tuple):
            # Si ya es una tupla, asegurarse de que esté en formato correcto
            if all(isinstance(c, (int, float)) for c in color):
                if all(0 <= c <= 1 for c in color):
                    # Ya está en formato 0-1
                    return color
                elif all(0 <= c <= 255 for c in color):
                    # Convertir de 0-255 a 0-1
                    return (color[0]/255, color[1]/255, color[2]/255)
        return (0, 0, 0)  # Negro por defecto

    def save_as_pdf(self, output_path: str):
        try:
            self.doc.save(output_path, garbage=4, deflate=True, clean=True)
        finally:
            self.doc.close()
            # Limpiar archivos temporales (aunque PDF no crea muchos)
            self._cleanup_temp_files()

    def get_preview_pixmap(self, data_map: dict, font_map: dict):
        temp_doc = fitz.open(self.template_path)
        try:
            page = temp_doc.load_page(0)
            
            # Buscar todos los placeholders para la previsualización
            all_instances = {}
            for placeholder in data_map.keys():
                instances = page.search_for(placeholder)
                if instances:
                    all_instances[placeholder] = instances
            
            for placeholder, instances in all_instances.items():
                value = data_map[placeholder]
                font_info = font_map.get(placeholder, {'family': 'Arial', 'size': 12, 'bold': False, 'color': (0, 0, 0)})
                
                # Determinar alineación
                is_folio = placeholder == "{{FOLIO}}"
                align_center = not is_folio
                
                if instances:
                    inst = instances[0]
                    page.add_redact_annot(inst)
                    page.apply_redactions()
                    
                    x_insert, y_insert, final_size, font_name = self._get_text_fit_info(
                        inst, value, font_info, align_center=align_center
                    )
                    
                    color = self._parse_color(font_info.get('color', (0, 0, 0)))
                    
                    page.insert_text((x_insert, y_insert), value, fontsize=final_size, fontname=font_name, color=color)

            pix = page.get_pixmap()
            return pix
        finally:
            temp_doc.close()

    def _get_pdf_font(self, family, bold):
        family_lower = family.lower()
        if "arial" in family_lower or "helvetica" in family_lower: 
            return "helv-bold" if bold else "helv"
        if "times" in family_lower: 
            return "tibo" if bold else "timo"
        if "courier" in family_lower: 
            return "cobo" if bold else "cour"
        return "helv-bold" if bold else "helv"

class OfficeProcessor(BaseProcessor):
    def _convert_to_pdf_with_com(self, input_path: str, output_path: str, app_name: str, format_type: int):
        app = None
        doc = None
        max_retries = 3
        retry_delay = 2
        
        for attempt in range(max_retries):
            try:
                app = comtypes.client.CreateObject(app_name)
                app.Visible = False
                
                if "powerpoint" in app_name.lower():
                    doc = app.Presentations.Open(input_path)
                else:
                    doc = app.Documents.Open(input_path)
                    
                doc.SaveAs(output_path, FileFormat=format_type)
                break  # Éxito, salir del loop de reintentos
                
            except Exception as e:
                print(f"⚠️ Intento {attempt + 1} fallado: {e}")
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                else:
                    raise e
            finally:
                # Cerrar documentos y aplicación
                if doc:
                    try:
                        doc.Close(False)
                    except:
                        pass
                if app:
                    try:
                        app.Quit()
                    except:
                        pass

class DocxProcessor(OfficeProcessor):
    def process(self, data_map: dict, font_map: dict): 
        self.doc = Document(self.template_path)
        
        # Procesar todos los placeholders en el data_map
        for placeholder, value in data_map.items():
            font_info = font_map.get(placeholder, {'family': 'Arial', 'size': 12, 'bold': False, 'color': '#000000'})
            
            # Determinar alineación
            is_folio = placeholder == "{{FOLIO}}"
            
            for paragraph in self.doc.paragraphs:
                if placeholder in paragraph.text:
                    # Reemplazar el texto manteniendo el formato
                    self._replace_text_in_paragraph(paragraph, placeholder, str(value), font_info, is_folio)
            
            # Procesar también en tablas
            for table in self.doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            if placeholder in paragraph.text:
                                self._replace_text_in_paragraph(paragraph, placeholder, str(value), font_info, is_folio)

    def _replace_text_in_paragraph(self, paragraph, placeholder, value, font_info, is_folio):
        """Reemplaza texto en un párrafo manteniendo el formato"""
        if placeholder in paragraph.text:
            # Reemplazar todo el texto del párrafo
            original_text = paragraph.text
            new_text = original_text.replace(placeholder, value)
            
            # Limpiar el párrafo y agregar nuevo texto con formato
            paragraph.clear()
            run = paragraph.add_run(new_text)
            
            # Aplicar formato
            font = run.font
            font.name = font_info['family']
            font.size = Pt(font_info['size'])
            font.bold = font_info['bold']
            
            # Manejar color
            color = font_info.get('color', '#000000')
            if isinstance(color, str) and color.startswith('#'):
                # Convertir hex a RGB
                color = color.lstrip('#')
                r = int(color[0:2], 16)
                g = int(color[2:4], 16)
                b = int(color[4:6], 16)
                font.color.rgb = RGBColor(r, g, b)
            
            # Aplicar alineación diferenciada
            if is_folio:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Folio alineado a la izquierda
            else:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER  # TEXT_1 y TEXT_2 centrados

    def save_as_pdf(self, output_path: str):
        temp_dir = os.path.dirname(output_path)
        
        # Usar nombre temporal más seguro
        temp_docx_path = os.path.join(temp_dir, f"temp_docx_{os.getpid()}_{int(time.time())}.docx")
        self.temp_files.append(temp_docx_path)
        
        try:
            # 1. Guardar el DOCX modificado
            self.doc.save(temp_docx_path)
            
            # 2. Convertir a PDF
            self._convert_to_pdf_with_com(
                os.path.abspath(temp_docx_path),
                os.path.abspath(output_path),
                'Word.Application',
                17  # wdFormatPDF
            )
            
        finally:
            # LIMPIEZA MEJORADA - Intentar múltiples veces
            self._cleanup_with_retry()

    def _cleanup_with_retry(self, max_retries=5, retry_delay=1):
        """Limpia archivos temporales con reintentos"""
        for attempt in range(max_retries):
            try:
                self._cleanup_temp_files()
                break  # Éxito
            except Exception as e:
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                else:
                    print(f"❌ No se pudieron eliminar algunos archivos temporales: {e}")

class PptxProcessor(OfficeProcessor):
    def process(self, data_map: dict, font_map: dict): 
        self.doc = Presentation(self.template_path)
        
        for placeholder, value in data_map.items():
            font_info = font_map.get(placeholder, {'family': 'Arial', 'size': 18, 'bold': False, 'color': '#000000'})
            
            # Determinar alineación
            is_folio = placeholder == "{{FOLIO}}"
            
            for slide in self.doc.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    # Procesar texto en formas
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if placeholder in paragraph.text:
                                # Reemplazar texto
                                original_text = paragraph.text
                                new_text = original_text.replace(placeholder, str(value))
                                paragraph.text = new_text
                                
                                # Aplicar alineación
                                if is_folio:
                                    paragraph.alignment = 1  # LEFT align para folio
                                else:
                                    paragraph.alignment = 2  # CENTER align para TEXT_1 y TEXT_2
                                
                                # Aplicar formato si hay runs
                                if paragraph.runs:
                                    font = paragraph.runs[0].font
                                    font.name = font_info['family']
                                    font.size = PptxPt(font_info['size'])
                                    font.bold = font_info['bold']
                                    
                                    # Manejar color
                                    color = font_info.get('color', '#000000')
                                    if isinstance(color, str) and color.startswith('#'):
                                        color = color.lstrip('#')
                                        r = int(color[0:2], 16)
                                        g = int(color[2:4], 16)
                                        b = int(color[4:6], 16)
                                        font.color.rgb = PptxRGBColor(r, g, b)

    def save_as_pdf(self, output_path: str):
        temp_dir = os.path.dirname(output_path)
        temp_pptx_path = os.path.join(temp_dir, f"temp_pptx_{os.getpid()}_{int(time.time())}.pptx")
        self.temp_files.append(temp_pptx_path)

        try:
            self.doc.save(temp_pptx_path)
            self._convert_to_pdf_with_com(
                os.path.abspath(temp_pptx_path),
                os.path.abspath(output_path),
                'Powerpoint.Application',
                32  # ppSaveAsPDF
            )
        finally:
            # LIMPIEZA MEJORADA
            self._cleanup_with_retry()

    def _cleanup_with_retry(self, max_retries=5, retry_delay=1):
        """Limpia archivos temporales con reintentos"""
        for attempt in range(max_retries):
            try:
                self._cleanup_temp_files()
                break  # Éxito
            except Exception as e:
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                else:
                    print(f"❌ No se pudieron eliminar algunos archivos temporales: {e}")

def get_processor(template_path: str):
    ext = os.path.splitext(template_path)[1].lower()
    if ext == '.pdf':
        return PdfProcessor(template_path)
    elif ext == '.docx':
        return DocxProcessor(template_path)
    elif ext == '.pptx':
        return PptxProcessor(template_path)
    else:
        raise ValueError(f"Tipo de archivo no soportado: {ext}")
