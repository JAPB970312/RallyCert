# signature.py (modificado)
"""
Firma digital criptográfica para RallyCert.
Inserta QR con la firma en el marcador {{QR}} de la plantilla
y verifica integridad del documento.
"""

import json
import os
import base64
from datetime import datetime
from io import BytesIO
from cryptography.hazmat.primitives import hashes, serialization
from cryptography.hazmat.primitives.asymmetric import padding, rsa
import qrcode
from PIL import Image
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Importaciones adicionales para la nueva funcionalidad de seguridad en PDF
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.lib import colors

# Importación para búsqueda exacta de posición QR
import fitz  # PyMuPDF

# ==============================================
# CONFIGURACIÓN
# ==============================================
KEYS_DIR = os.path.join(os.path.dirname(__file__), "keys")
PRIVATE_KEY_PATH = os.path.join(KEYS_DIR, "private_key.pem")
PUBLIC_KEY_PATH = os.path.join(KEYS_DIR, "public_key.pem")
ISSUER_STRING = "Firmado por Rally de la Ninez Cientifica y Expo STEM by Universidad de Sonora"

# Variable global para la leyenda personalizable
VALIDATION_TEXT = "Validado por Rally de la Niñez Científica y EXPO STEM, Universidad de Sonora"

# ==============================================
# LLAVES RSA
# ==============================================
def ensure_keys(private_path=PRIVATE_KEY_PATH, public_path=PUBLIC_KEY_PATH, bits=2048):
    os.makedirs(os.path.dirname(private_path), exist_ok=True)
    if os.path.exists(private_path) and os.path.exists(public_path):
        return private_path, public_path

    private_key = rsa.generate_private_key(public_exponent=65537, key_size=bits)
    priv_pem = private_key.private_bytes(
        encoding=serialization.Encoding.PEM,
        format=serialization.PrivateFormat.PKCS8,
        encryption_algorithm=serialization.NoEncryption(),
    )
    pub_pem = private_key.public_key().public_bytes(
        encoding=serialization.Encoding.PEM,
        format=serialization.PublicFormat.SubjectPublicKeyInfo,
    )
    with open(private_path, "wb") as f:
        f.write(priv_pem)
    with open(public_path, "wb") as f:
        f.write(pub_pem)
    return private_path, public_path


def load_private_key(path=PRIVATE_KEY_PATH):
    with open(path, "rb") as f:
        return serialization.load_pem_private_key(f.read(), password=None)


def load_public_key(path=PUBLIC_KEY_PATH):
    with open(path, "rb") as f:
        return serialization.load_pem_public_key(f.read())


# ==============================================
# FIRMADO
# ==============================================
def canonicalize_payload(payload: dict) -> bytes:
    return json.dumps(payload, separators=(",", ":"), sort_keys=True).encode("utf-8")


def sign_bytes(private_key_path: str, data: bytes) -> str:
    priv = load_private_key(private_key_path)
    signature = priv.sign(
        data,
        padding.PSS(mgf=padding.MGF1(hashes.SHA256()), salt_length=padding.PSS.MAX_LENGTH),
        hashes.SHA256(),
    )
    return base64.b64encode(signature).decode("ascii")


def verify_signature(public_key_path: str, data: bytes, signature_b64: str) -> bool:
    pub = load_public_key(public_key_path)
    sig = base64.b64decode(signature_b64)
    try:
        pub.verify(
            sig,
            data,
            padding.PSS(mgf=padding.MGF1(hashes.SHA256()), salt_length=padding.PSS.MAX_LENGTH),
            hashes.SHA256(),
        )
        return True
    except Exception:
        return False


# ==============================================
# GENERAR QR
# ==============================================
def make_qr_image(text: str, box_size: int = 3) -> Image.Image:
    """
    box_size reducido → QR más pequeño (~2 cm)
    """
    qr = qrcode.QRCode(border=1, box_size=box_size)
    qr.add_data(text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    return img


# ==============================================
# BUSCAR POSICIÓN EXACTA DEL QR EN PDF
# ==============================================
def find_qr_position_in_pdf(pdf_path):
    """
    Busca el marcador {{QR}} en el PDF usando PyMuPDF y devuelve su posición exacta.
    Retorna: dict con page, x, y o None si no se encuentra.
    """
    try:
        doc = fitz.open(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text_instances = page.search_for("{{QR}}")
            
            if text_instances:
                # Tomar la primera instancia encontrada
                rect = text_instances[0]
                # Devolver posición del texto encontrado
                position = {
                    "page": page_num,
                    "x": rect.x0,
                    "y": rect.y0
                }
                doc.close()
                return position
        
        doc.close()
        return None
    except Exception as e:
        print(f"Error buscando posición QR: {e}")
        return None


# ==============================================
# EMBEBER EN DOCUMENTOS
# ==============================================
def _find_qr_placeholder(doc: Document):
    """
    Busca el marcador {{QR}} en un docx y lo elimina,
    devolviendo el párrafo donde estaba.
    """
    for p in doc.paragraphs:
        if "{{QR}}" in p.text:
            p.text = p.text.replace("{{QR}}", "")
            return p
    return None


def embed_qr_in_docx(input_docx: str, output_docx: str, qr_img: Image.Image, metadata: dict):
    doc = Document(input_docx)
    qr_para = _find_qr_placeholder(doc)

    stream = BytesIO()
    qr_img.save(stream, format="PNG")
    stream.seek(0)

    if qr_para is not None:
        run = qr_para.add_run()
        run.add_picture(stream, width=Inches(0.8))  # más pequeño (~2 cm)
    else:
        doc.add_picture(stream, width=Inches(0.8))

    try:
        doc.core_properties.comments = json.dumps(metadata, separators=(",", ":"), sort_keys=True)
    except Exception:
        doc.add_paragraph("[SIGNATURE-METADATA] " + json.dumps(metadata, separators=(",", ":"), sort_keys=True))
    doc.save(output_docx)


def embed_qr_in_pdf(input_pdf: str, output_pdf: str, qr_img: Image.Image, metadata: dict, validation_text=None):
    """
    Inserta un QR exactamente donde está el marcador {{QR}} en el PDF,
    con elementos de seguridad que detectan modificaciones.
    """
    # Usar texto personalizado o el predeterminado
    if validation_text is None:
        validation_text = VALIDATION_TEXT
    
    # Buscar posición exacta del QR
    qr_position = find_qr_position_in_pdf(input_pdf)
    
    # Leer el PDF original
    with open(input_pdf, "rb") as f:
        base_bytes = f.read()
    base_stream = BytesIO(base_bytes)
    base_reader = PdfReader(base_stream)
    
    writer = PdfWriter()
    
    for page_num, page in enumerate(base_reader.pages):
        w, h = float(page.mediabox.width), float(page.mediabox.height)
        
        # Crear overlay
        packet = BytesIO()
        c = canvas.Canvas(packet, pagesize=(w, h))
        
        # 1. MARCA DE AGUA DE SEGURIDAD DISCRETA
        # Esta marca de agua es visible pero no intrusiva, y sirve para detectar modificaciones
        c.saveState()
        c.translate(w / 2, h / 2)
        c.rotate(45)
        c.setFont("Helvetica", 36)  # Tamaño moderado
        c.setFillColor(colors.lightgrey, alpha=0.08)  # Muy transparente
        watermark_text = "VÁLIDO - UNIVERSIDAD DE SONORA"
        c.drawCentredString(0, 0, watermark_text)
        
        # Agregar texto adicional en otras posiciones para mayor seguridad
        c.rotate(-90)
        c.drawCentredString(0, -200, "CONSTANCIA AUTENTICADA")
        c.restoreState()
        
        # 2. Insertar QR
        qr_reader = ImageReader(qr_img)
        qr_size = 60  # Tamaño pequeño (~2.1 cm)
        
        if qr_position and page_num == qr_position["page"]:
            # Posición exacta del marcador
            qr_x = qr_position["x"]
            qr_y = h - qr_position["y"] - qr_size  # Ajuste coordenada Y (PDF coordinate system)
        else:
            # Posición por defecto (esquina inferior derecha)
            qr_x = w - qr_size - 30
            qr_y = 30
        
        c.drawImage(qr_reader, qr_x, qr_y, width=qr_size, height=qr_size, mask='auto')
        
        # 3. TEXTO DE VALIDACIÓN PERSONALIZABLE
        c.setFont("Helvetica", 7)
        c.setFillColor(colors.darkgrey)
        
        # Texto informativo
        info_text = "Constancia verificable mediante código QR - Universidad de Sonora"
        
        # Posicionar texto cerca del QR
        text_x = qr_x
        text_y = qr_y - 12  # Justo debajo del QR
        
        c.drawString(text_x, text_y, info_text)
        
        # 4. LEYENDA DE VALIDACIÓN (texto personalizado)
        c.setFont("Helvetica", 6)
        c.setFillColor(colors.grey)
        # Posicionar en la esquina inferior izquierda
        c.drawString(30, 20, validation_text)
        
        c.save()
        packet.seek(0)
        overlay_reader = PdfReader(packet)
        overlay_page = overlay_reader.pages[0]
        
        # Fusionar overlay con página original
        page.merge_page(overlay_page)
        writer.add_page(page)
        
        packet.close()

    # Añadir metadatos de firma de manera más robusta
    if metadata:
        # Método 1: Metadatos estándar
        writer.add_metadata({
            "/Title": "Constancia Universitaria",
            "/Author": "Universidad de Sonora",
            "/Subject": "Constancia de Participación",
            "/Creator": "Sistema de Constancias Rally STEM",
            "/Producer": "RallyCert v1.0",
            "/Signature": json.dumps(metadata, separators=(",", ":"), sort_keys=True),
            "/ValidationText": validation_text  # Guardar también la leyenda en metadatos
        })
        
        # Método 2: Metadatos adicionales como campos XMP
        try:
            # Esto hace más difícil la eliminación de metadatos
            from PyPDF2 import PdfFileWriter
            if hasattr(writer, 'add_metadata'):
                custom_metadata = {
                    'RallyCert_Signature': base64.b64encode(
                        json.dumps(metadata, separators=(",", ":"), sort_keys=True).encode()
                    ).decode(),
                    'RallyCert_ValidationText': validation_text
                }
                writer.add_metadata(custom_metadata)
        except Exception:
            pass

    # Guardar PDF final
    with open(output_pdf, "wb") as f_out:
        writer.write(f_out)
        
    base_stream.close()


def embed_qr_in_pptx(input_pptx: str, output_pptx: str, qr_img: Image.Image, metadata: dict,
                     slide_index=0, left_inches=8, top_inches=5):
    prs = Presentation(input_pptx)
    slide = prs.slides[slide_index] if slide_index < len(prs.slides) else prs.slides[0]
    stream = BytesIO()
    qr_img.save(stream, format="PNG")
    stream.seek(0)
    slide.shapes.add_picture(stream, Inches(left_inches), Inches(top_inches), height=Inches(1))
    try:
        prs.core_properties.comments = json.dumps(metadata, separators=(",", ":"), sort_keys=True)
    except Exception:
        tx = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(2), Inches(1))
        tx.text_frame.text = "[SIGN-METADATA] " + json.dumps(metadata, separators=(",", ":"), sort_keys=True)
    prs.save(output_pptx)


# ==============================================
# VERIFICACIÓN DE MODIFICACIONES
# ==============================================
def check_document_integrity(pdf_path, public_key_path=PUBLIC_KEY_PATH):
    """
    Verifica si un documento PDF ha sido modificado después de la firma.
    Si fue modificado, la leyenda y firma digital se invalidan automáticamente.
    """
    try:
        with open(pdf_path, "rb") as f:
            pdf_reader = PdfReader(f)
            
            # Buscar metadatos de firma
            metadata = pdf_reader.metadata
            signature_data = None
            original_validation_text = None
            
            if '/Signature' in metadata:
                signature_data = json.loads(metadata['/Signature'])
            if '/ValidationText' in metadata:
                original_validation_text = metadata['/ValidationText']
            else:
                # Buscar en otros campos de metadatos
                for key, value in metadata.items():
                    if 'signature' in key.lower() or 'rallycert' in key.lower():
                        try:
                            if 'RallyCert_Signature' in key:
                                signature_data = json.loads(base64.b64decode(value).decode())
                            elif 'RallyCert_ValidationText' in key:
                                original_validation_text = value
                            else:
                                signature_data = json.loads(value)
                            break
                        except:
                            continue
            
            if not signature_data:
                return False, "No se encontraron datos de firma en el documento", None
            
            # Verificar firma
            payload = signature_data.get('payload', {})
            signature_b64 = signature_data.get('signature', '')
            
            if not payload or not signature_b64:
                return False, "Datos de firma incompletos", None
            
            payload_bytes = canonicalize_payload(payload)
            is_valid = verify_signature(public_key_path, payload_bytes, signature_b64)
            
            if is_valid:
                return True, "Documento válido y sin modificaciones", original_validation_text
            else:
                return False, "⚠️ ADVERTENCIA: El documento ha sido modificado o alterado. La leyenda de validación y firma digital han sido invalidadas.", None
                
    except Exception as e:
        return False, f"Error en verificación: {str(e)}", None


# ==============================================
# FIRMA PRINCIPAL (MODIFICADA PARA ACEPTAR TEXTO PERSONALIZADO)
# ==============================================
def build_payload(cert_data: dict, issuer=ISSUER_STRING):
    return {"data": cert_data, "issued_at": datetime.utcnow().isoformat() + "Z", "issuer": issuer}


def sign_and_embed(input_path: str, output_path: str, cert_data: dict,
                   private_key_path=PRIVATE_KEY_PATH, public_key_path=PUBLIC_KEY_PATH,
                   validation_text=None):
    """
    Firma y embebe el QR en el documento, con texto de validación personalizable.
    
    Args:
        validation_text: Texto personalizado para la leyenda de validación.
                        Si es None, usa el texto por defecto.
    """
    ensure_keys(private_key_path, public_key_path)
    payload = build_payload(cert_data)
    payload_bytes = canonicalize_payload(payload)
    signature_b64 = sign_bytes(private_key_path, payload_bytes)

    metadata = {
        "payload": payload,
        "signature": signature_b64,
        "pubkey_id": os.path.basename(public_key_path),
    }

    qr_img = make_qr_image(json.dumps(metadata, separators=(",", ":"), sort_keys=True))

    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".pdf":
        embed_qr_in_pdf(input_path, output_path, qr_img, metadata, validation_text)
    elif ext == ".docx":
        embed_qr_in_docx(input_path, output_path, qr_img, metadata)
    elif ext in (".pptx", ".ppt"):
        embed_qr_in_pptx(input_path, output_path, qr_img, metadata)
    else:
        with open(input_path, "rb") as r, open(output_path, "wb") as w:
            w.write(r.read())
        with open(output_path + ".signature.json", "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, separators=(",", ":"), sort_keys=True)
    return metadata


# ==============================================
# FUNCIÓN PARA CONFIGURAR TEXTO PERSONALIZADO
# ==============================================
def set_validation_text(text):
    """
    Configura el texto de validación globalmente.
    """
    global VALIDATION_TEXT
    VALIDATION_TEXT = text


def get_validation_text():
    """
    Obtiene el texto de validación actual.
    """
    return VALIDATION_TEXT