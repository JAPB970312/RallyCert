# verify_signature.py
"""
Verificador de firmas digitales de constancias RallyCert.
Detecta modificaciones y marca el documento como MODIFICADO.
"""

import json
import os
from io import BytesIO
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from signature import verify_signature, canonicalize_payload, PUBLIC_KEY_PATH


def _extract_from_pdf(path):
    reader = PdfReader(path)
    info = reader.metadata or {}
    sig = info.get("/Signature")
    if sig:
        try:
            return json.loads(sig)
        except Exception:
            pass
    return None


def _extract_from_docx(path):
    doc = Document(path)
    c = doc.core_properties.comments
    if c:
        try:
            return json.loads(c)
        except Exception:
            pass
    for p in doc.paragraphs:
        if "[SIGNATURE-METADATA]" in p.text:
            try:
                return json.loads(p.text.split("]")[-1].strip())
            except Exception:
                pass
    return None


def _extract_from_pptx(path):
    prs = Presentation(path)
    c = prs.core_properties.comments
    if c:
        try:
            return json.loads(c)
        except Exception:
            pass
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and "[SIGN-METADATA]" in shape.text:
                try:
                    return json.loads(shape.text.split("]")[-1].strip())
                except Exception:
                    pass
    return None


def extract_metadata(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _extract_from_pdf(path)
    if ext == ".docx":
        return _extract_from_docx(path)
    if ext in (".pptx", ".ppt"):
        return _extract_from_pptx(path)
    side = path + ".signature.json"
    if os.path.exists(side):
        with open(side, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def mark_pdf_modified(input_pdf, output_pdf):
    r = PdfReader(input_pdf)
    w = PdfWriter()
    for p in r.pages:
        w.add_page(p)
    packet = BytesIO()
    c = canvas.Canvas(packet, pagesize=letter)
    c.setFont("Helvetica-Bold", 48)
    c.setFillAlpha(0.3)
    c.drawString(100, 400, "MODIFICADO / MANIPULADO")
    c.save()
    packet.seek(0)
    overlay = PdfReader(packet)
    try:
        w.pages[0].merge_page(overlay.pages[0])
    except Exception:
        pass
    with open(output_pdf, "wb") as f:
        w.write(f)


def mark_docx_modified(input_docx, output_docx):
    d = Document(input_docx)
    d.add_paragraph("*** MODIFICADO / MANIPULADO ***")
    d.save(output_docx)


def mark_pptx_modified(input_pptx, output_pptx):
    prs = Presentation(input_pptx)
    for s in prs.slides:
        for sh in s.shapes:
            if hasattr(sh, "text_frame"):
                sh.text_frame.text = "[MODIFICADO] " + sh.text_frame.text
    prs.save(output_pptx)


def verify_document(path, public_key_path=PUBLIC_KEY_PATH):
    meta = extract_metadata(path)
    if not meta:
        return False, "No se encontró firma digital", None
    payload = meta.get("payload")
    sig = meta.get("signature")
    if not payload or not sig:
        return False, "Firma dañada o incompleta", meta

    ok = verify_signature(public_key_path, canonicalize_payload(payload), sig)
    if ok:
        return True, "✅ Firma válida e íntegra", meta

    ext = os.path.splitext(path)[1].lower()
    out = os.path.splitext(path)[0] + "_MODIFICADO" + ext
    if ext == ".pdf":
        mark_pdf_modified(path, out)
    elif ext == ".docx":
        mark_docx_modified(path, out)
    elif ext in (".pptx", ".ppt"):
        mark_pptx_modified(path, out)
    return False, f"❌ Documento manipulado, se generó {out}", meta
